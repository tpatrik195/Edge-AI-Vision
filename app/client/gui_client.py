import tkinter as tk
from tkinter import filedialog, messagebox
import cv2
import requests
import numpy as np
import threading
from PIL import Image, ImageTk
import os
from app.models.segmentation import BackgroundSegmenter
from pdf2image import convert_from_path
from pptx import Presentation
from fastapi import FastAPI, Request
import uvicorn

SERVER_URL = "http://127.0.0.1:8000"
WEBHOOK_URL = "http://127.0.0.1:9001/webhook"

background_segmenter = None
scaling_factor = 1.0

root = tk.Tk()
root.title("Edge AI Vision")
root.geometry("800x600")

current_image = None
pdf_images = []
ppt_slides = []
current_page = 0
ppt_current_slide = 0
cap = None
camera_active = False
webhook_subscribed = False

def display_image(image):
    global current_image
    if image is not None:
        img = Image.fromarray(image)
        img = img.resize((600, 400))
        img_tk = ImageTk.PhotoImage(img)

        if current_image is None:
            current_image = tk.Label(root, image=img_tk)
            current_image.image = img_tk
            current_image.pack(pady=20)
        else:
            current_image.configure(image=img_tk)
            current_image.image = img_tk

def convert_to_image(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.pdf':
        images = convert_from_path(file_path)
        return [np.array(image) for image in images]
    
    elif ext == '.pptx':
        prs = Presentation(file_path)
        slides = []
        for slide in prs.slides:
            if slide.shapes:
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        img = np.array(shape.image)
                        slides.append(img)
        return slides
    
    return None

def open_file():
    global background_segmenter, pdf_images, ppt_slides, current_page, ppt_current_slide
    file_path = filedialog.askopenfilename(filetypes=[("Images", "*.jpg;*.png"), ("PDF", "*.pdf"), ("PPT", "*.pptx")])
    
    if file_path:
        ext = os.path.splitext(file_path)[1].lower()

        if ext == '.pdf' or ext == '.pptx':
            images = convert_to_image(file_path)
            if ext == '.pdf':
                pdf_images = images
                current_page = 0
            elif ext == '.pptx':
                ppt_slides = images
                ppt_current_slide = 0
            image = images[0]
        else:
            image = cv2.imread(file_path)
        
        if image is not None:
            background_segmenter = BackgroundSegmenter(image)
            display_image(image)
        # else:
            # messagebox.showerror("Error", "Failed to load image from file.")

def next_page(event=None):
    global current_page, ppt_current_slide, pdf_images, ppt_slides, background_segmenter

    if pdf_images:
        if current_page < len(pdf_images) - 1:
            current_page += 1
            image = pdf_images[current_page]
            background_segmenter = BackgroundSegmenter(image)
            segmented_image = background_segmenter.segment_background(image)
            display_image(segmented_image)
    elif ppt_slides:
        if ppt_current_slide < len(ppt_slides) - 1:
            ppt_current_slide += 1
            image = ppt_slides[ppt_current_slide]
            background_segmenter = BackgroundSegmenter(image)
            segmented_image = background_segmenter.segment_background(image)
            display_image(segmented_image)

def previous_page(event=None):
    global current_page, ppt_current_slide, pdf_images, ppt_slides, background_segmenter

    if pdf_images:
        if current_page > 0:
            current_page -= 1
            image = pdf_images[current_page]
            background_segmenter = BackgroundSegmenter(image)
            segmented_image = background_segmenter.segment_background(image)
            display_image(segmented_image)
    elif ppt_slides:
        if ppt_current_slide > 0:
            ppt_current_slide -= 1
            image = ppt_slides[ppt_current_slide]
            background_segmenter = BackgroundSegmenter(image)
            segmented_image = background_segmenter.segment_background(image)
            display_image(segmented_image)

def subscribe_webhook():
    global webhook_subscribed
    if webhook_subscribed:
        # messagebox.showinfo("Info", "Already subscribed.")
        return
    
    response = requests.post(f"{SERVER_URL}/subscribe_webhook", params={"url": WEBHOOK_URL})
    if response.status_code == 200:
        webhook_subscribed = True
        # messagebox.showinfo("Info", "Subscribed to webhook.")
        start_webcam()
    # else:
        # messagebox.showerror("Error", "Failed to subscribe to webhook.")

def unsubscribe_webhook():
    global webhook_subscribed
    if not webhook_subscribed:
        # messagebox.showinfo("Info", "Not subscribed.")
        return
    
    response = requests.post(f"{SERVER_URL}/unsubscribe_webhook", params={"url": WEBHOOK_URL})
    if response.status_code == 200:
        webhook_subscribed = False
        # messagebox.showinfo("Info", "Unsubscribed from webhook.")
        stop_webcam()
    # else:
        # messagebox.showerror("Error", "Failed to unsubscribe from webhook.")

def start_webcam():
    global cap, camera_active
    if camera_active:
        return
    
    cap = cv2.VideoCapture(0)
    if not cap.isOpened():
        # messagebox.showerror("Error", "Cannot open webcam!")
        return
    
    camera_active = True
    threading.Thread(target=send_frames, daemon=True).start()

def stop_webcam():
    global cap, camera_active
    if not camera_active:
        return
    
    camera_active = False
    cap.release()
    cv2.destroyAllWindows()

def send_frames():
    global cap
    while camera_active:
        ret, frame = cap.read()
        if not ret:
            print("Error: Cannot read frame!")
            break
        
        _, img_encoded = cv2.imencode('.jpg', frame)
        requests.post(f"{SERVER_URL}/process_frame", files={"frame": img_encoded.tobytes()})

        if background_segmenter is not None:
            segmented_frame = background_segmenter.segment_background(frame)
            if segmented_frame is not None:
                display_image(segmented_frame)

        if cv2.waitKey(1) == 27:
            stop_webcam()
            break

app = FastAPI()

@app.post("/webhook")
async def webhook(request: Request):
    data = await request.json()
    print(f"Gesture Event Received: {data}")
    received_gesture = data.get("gesture", "unknown")

    if received_gesture == "swipe right":
        root.after(0, next_page)
    elif received_gesture == "swipe left":
        root.after(0, previous_page)
    # elif received_gesture == "+":
    #         root.after(0, lambda: resize_segmented_image(1.1))
    # elif received_gesture == "-":
    #         root.after(0, lambda: resize_segmented_image(0.9))

    return {"status": "received"}

def run_fastapi_server():
    uvicorn.run(app, host="127.0.0.1", port=9001)

threading.Thread(target=run_fastapi_server, daemon=True).start()

root.bind("<Left>", previous_page)
root.bind("<Right>", next_page)

open_button = tk.Button(root, text="Open File", command=open_file)
open_button.pack(pady=10)

subscribe_button = tk.Button(root, text="Subscribe", command=subscribe_webhook)
subscribe_button.pack(pady=10)

unsubscribe_button = tk.Button(root, text="Unsubscribe", command=unsubscribe_webhook)
unsubscribe_button.pack(pady=10)

root.mainloop()
